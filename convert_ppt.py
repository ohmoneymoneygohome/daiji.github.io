import pptx
import os

def extract_ppt_content(ppt_path):
    """提取PPT内容"""
    try:
        presentation = pptx.Presentation(ppt_path)
        content = []
        
        for i, slide in enumerate(presentation.slides):
            slide_content = {
                'slide_num': i + 1,
                'title': '',
                'content': []
            }
            
            # 提取标题
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    if shape.text and not slide_content['title']:
                        slide_content['title'] = shape.text
                    elif shape.text:
                        slide_content['content'].append(shape.text)
            
            content.append(slide_content)
        
        return content
    except Exception as e:
        print(f"Error extracting PPT content: {e}")
        return []

def generate_html(content, output_path):
    """生成HTML文件"""
    html_template = """
<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>代际：时间的社会学与文明的辩证法</title>
    <script src="https://cdn.tailwindcss.com"></script>
    <link href="https://cdn.jsdelivr.net/npm/font-awesome@4.7.0/css/font-awesome.min.css" rel="stylesheet">
    <script>
        tailwind.config = {
            theme: {
                extend: {
                    colors: {
                        primary: '#4F46E5',
                        secondary: '#8B5CF6',
                    },
                    fontFamily: {
                        sans: ['Inter', 'system-ui', 'sans-serif'],
                    },
                }
            }
        }
    </script>
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');
        
        body {
            font-family: 'Inter', sans-serif;
        }
        
        .slide {
            scroll-margin-top: 80px;
        }
        
        .fade-in {
            animation: fadeIn 0.5s ease-in-out;
        }
        
        @keyframes fadeIn {
            from {
                opacity: 0;
                transform: translateY(20px);
            }
            to {
                opacity: 1;
                transform: translateY(0);
            }
        }
    </style>
</head>
<body class="bg-gray-50">
    <!-- 导航栏 -->
    <nav class="fixed w-full top-0 z-50 bg-white shadow-md">
        <div class="container mx-auto px-4 sm:px-6 lg:px-8">
            <div class="flex justify-between items-center h-16">
                <div class="flex items-center">
                    <i class="fa fa-book text-primary text-2xl mr-2"></i>
                    <span class="text-xl font-bold text-gray-900">代际：时间的社会学与文明的辩证法</span>
                </div>
                <div class="hidden md:flex space-x-6">
                    <a href="index.html" class="text-gray-700 hover:text-primary transition-colors">首页</a>
                    <a href="#intro" class="text-gray-700 hover:text-primary transition-colors">介绍</a>
                </div>
            </div>
        </div>
    </nav>

    <!-- 主内容 -->
    <main class="container mx-auto px-4 sm:px-6 lg:px-8 pt-24 pb-16">
        <!-- 标题部分 -->
        <section id="intro" class="mb-16 text-center">
            <h1 class="text-4xl sm:text-5xl font-bold text-gray-900 mb-6">代际：时间的社会学与文明的辩证法</h1>
            <p class="text-xl text-gray-600 max-w-3xl mx-auto">探索代际之间的时间社会学与文明发展的辩证关系</p>
        </section>

        <!-- PPT内容部分 -->
        <div class="space-y-16">
            {{SLIDES}}
        </div>
    </main>

    <!-- 页脚 -->
    <footer class="bg-gray-900 text-white py-12">
        <div class="container mx-auto px-4 sm:px-6 lg:px-8">
            <div class="text-center">
                <p class="text-gray-400">代际：时间的社会学与文明的辩证法</p>
                <p class="text-gray-500 text-sm mt-2">© 2026 代际知识库</p>
            </div>
        </div>
    </footer>

    <script>
        // 平滑滚动
        document.querySelectorAll('a[href^="#"]').forEach(anchor => {
            anchor.addEventListener('click', function (e) {
                e.preventDefault();
                document.querySelector(this.getAttribute('href')).scrollIntoView({
                    behavior: 'smooth'
                });
            });
        });
    </script>
</body>
</html>
    """
    
    # 生成幻灯片内容
    slides_html = ""
    for slide in content:
        slide_html = f"""
        <section class="slide bg-white rounded-xl shadow-md p-8 fade-in">
            <h2 class="text-2xl font-bold text-gray-900 mb-6">{slide['title']}</h2>
            <div class="space-y-4">
                {''.join([f'<p class="text-gray-700">{item}</p>' for item in slide['content']])}
            </div>
        </section>
        """
        slides_html += slide_html
    
    # 替换模板中的幻灯片内容
    html_content = html_template.replace('{{SLIDES}}', slides_html)
    
    # 写入HTML文件
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"HTML file generated at: {output_path}")

if __name__ == "__main__":
    ppt_path = "F:\\项目\\代际研究\\代际：时间的社会学与文明的辩证法.pptx"
    output_path = "F:\\daiji.github.io\\dialectic.html"
    
    content = extract_ppt_content(ppt_path)
    if content:
        generate_html(content, output_path)
    else:
        print("No content extracted from PPT")
